"""Contrato inicial del Informe Mensual de las 9 Atenciones Priorizadas."""
from pathlib import Path
from io import BytesIO
import sys
import tempfile

ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "backend"))

from modules.dbapi_compat import sqlite3  # noqa: E402
from modules.reportes_gerenciales.services import ReportesGerencialesService  # noqa: E402
from modules.reportes_gerenciales.atenciones_priorizadas import (  # noqa: E402
    ATENCIONES,
    AtencionesPriorizadasService,
)
from pptx import Presentation
from werkzeug.datastructures import FileStorage


def require(value, message):
    if not value:
        raise AssertionError(message)


with tempfile.TemporaryDirectory() as temp:
    db = Path(temp) / "test.sqlite3"
    conn = sqlite3.connect(str(db))
    conn.executescript("""
    CREATE TABLE fundaciones(id INTEGER PRIMARY KEY, nombre TEXT);
    CREATE TABLE master_ninos(
      id INTEGER PRIMARY KEY, fundacion_id INTEGER, activo INTEGER,
      unidad_servicio TEXT, documento TEXT, tipo_documento TEXT,
      tipo_beneficiario TEXT, edad_meses INTEGER, estado TEXT, sexo TEXT
    );
    CREATE TABLE master_talento_humano(
      id INTEGER PRIMARY KEY, fundacion_id INTEGER, activo INTEGER,
      nombre_completo TEXT, unidad_servicio TEXT, cargo TEXT
    );
    """)
    conn.execute("INSERT INTO fundaciones VALUES(1,'Fundación A')")
    conn.execute("INSERT INTO fundaciones VALUES(2,'Fundación B')")
    people = [
        (1, 1, 1, "UCA 1", "A1", "RC", "NIÑO", 4, "ACTIVO", "M"),
        (2, 1, 1, "UCA 1", "A2", "TI", "NIÑO", 10, "ACTIVO", "F"),
        (3, 1, 1, "UCA 1", "A3", "CC", "GESTANTE", 0, "ACTIVO", "F"),
        (4, 2, 1, "UCA B", "B1", "RC", "NIÑO", 24, "ACTIVO", "M"),
    ]
    conn.executemany("INSERT INTO master_ninos VALUES(?,?,?,?,?,?,?,?,?,?)", people)
    conn.execute("INSERT INTO master_talento_humano VALUES(1,1,1,'Docente A','UCA 1','DOCENTE')")
    conn.commit(); conn.close()

    base = ReportesGerencialesService(str(db), str(Path(temp) / "out"))
    base.init_schema()
    service = AtencionesPriorizadasService(base)
    require(len(ATENCIONES) == 9, "El catálogo no contiene nueve atenciones")

    data = service.consolidar(1, 8, 2026, 3)
    require(data["beneficiarios"] == 3, "Mezcló beneficiarios de otra fundación")
    require(data["rpp_total"] == 3, "Los rangos RPP no consolidan la población")
    doc = next(x for x in data["resultados"] if x["codigo"] == "DOCUMENTO_IDENTIDAD")
    require(doc["datos"]["tipos"] == {"CC": 1, "CE": 0, "RC": 1, "SD": 0, "TI": 1}, "Conteo documental incorrecto")
    require(abs(sum(doc["datos"]["tipos"].values()) - 3) == 0, "Documentos no suman beneficiarios")

    draft = service.guardar_borrador({
        "mes": 8, "anio": 2026, "contrato": "27005522025",
        "cobertura_contratada": 3, "modalidad": "Educación Inicial Propia Diaria",
    }, {"id": 1, "fundacion_id": 1})
    require(draft["id"] > 0 and draft["estado"] == "BORRADOR", "No guardó el borrador")
    conn = sqlite3.connect(str(db)); conn.row_factory = sqlite3.Row
    require(conn.execute("SELECT COUNT(*) c FROM rg9_resultados WHERE informe_id=?", (draft["id"],)).fetchone()["c"] == 9, "No guardó nueve resultados")
    require(conn.execute("SELECT COUNT(*) c FROM rg9_snapshots").fetchone()["c"] == 0, "Congeló un borrador")
    conn.close()

    try:
        service.aprobar(draft["id"], {"id": 1, "fundacion_id": 2})
        raise AssertionError("Fundación B aprobó informe de Fundación A")
    except LookupError:
        pass
    try:
        service.aprobar(draft["id"], {"id": 1, "fundacion_id": 1})
        raise AssertionError("Aprobó un informe con atenciones pendientes")
    except ValueError:
        pass

    # Completar las atenciones semiautomáticas permite aprobar, congelar y
    # producir los cuatro entregables sin recalcular el snapshot posteriormente.
    conn = sqlite3.connect(str(db))
    conn.execute("UPDATE rg9_resultados SET numerador=denominador,porcentaje=100,estado='COMPLETO' WHERE informe_id=?", (draft["id"],))
    conn.execute("UPDATE rg9_hallazgos SET estado='CERRADO' WHERE informe_id=?", (draft["id"],))
    conn.commit(); conn.close()
    approved = service.aprobar(draft["id"], {"id": 1, "fundacion_id": 1})
    require(len(approved["snapshot_hash"]) == 64, "No generó hash del snapshot")
    template_bytes=BytesIO(); template_prs=Presentation(); template_prs.slides.add_slide(template_prs.slide_layouts[0]); template_prs.save(template_bytes); template_bytes.seek(0)
    template=service.guardar_plantilla_pptx(FileStorage(stream=template_bytes,filename="plantilla_oficial_icbf.pptx"),{"version":"2026","fecha_vigencia":"2026-08-01"},{"id":1,"fundacion_id":1})
    require(template["estado"]=="ACTIVA" and template["diapositivas"]==1,"No activó la plantilla PPTX")
    require(service.plantilla_pptx_activa(2) is None,"La plantilla se cruzó con otra fundación")
    exports = service.generar_exportaciones(draft["id"], 1)
    require(set(exports) == {"pptx", "pdf", "xlsx", "zip"}, "Faltan exportaciones")
    require(all(Path(path).is_file() and Path(path).stat().st_size > 0 for path in exports.values()), "Alguna exportación está vacía")
    require(len(Presentation(exports["pptx"]).slides)>=11,"No conservó la plantilla ni agregó el informe")

print("REPORTES_9_ATENCIONES_V2_7_4_PASS")
