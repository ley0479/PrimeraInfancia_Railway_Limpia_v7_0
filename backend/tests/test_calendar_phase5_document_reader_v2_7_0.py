#!/usr/bin/env python3
"""Fase 5: lectores independientes, confianza y aprobación humana."""
from __future__ import annotations
import sys, tempfile
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

ROOT=Path(__file__).resolve().parents[2]; BACKEND=ROOT/'backend'; sys.path.insert(0,str(BACKEND))
from modules.calendario_inteligente.repository import CalendarioInteligenteRepository
from modules.calendario_inteligente.services import construir_preview_cronograma
from modules.seguridad.tenant_context import tenant_context

def req(ok,msg):
    if not ok: raise AssertionError(msg)

def run():
    with tempfile.TemporaryDirectory(prefix='pi-phase5-') as td:
        root=Path(td); headers=['Fecha','Actividad','TH a cargo']; values=['20/08/2026','Encuentro familiar','Psicosocial']
        xlsx=root/'planeacion.xlsx'; wb=Workbook(); ws=wb.active; ws.append(headers); ws.append(values); wb.save(xlsx)
        docx=root/'planeacion.docx'; doc=Document(); table=doc.add_table(rows=2,cols=3)
        for i,v in enumerate(headers): table.rows[0].cells[i].text=v
        for i,v in enumerate(values): table.rows[1].cells[i].text=v
        doc.save(docx)
        pptx=root/'planeacion.pptx'; prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[5]); box=slide.shapes.add_textbox(0,0,6000000,2000000); box.text_frame.text='Fecha | Actividad | TH a cargo\n20/08/2026 | Encuentro familiar | Psicosocial'; prs.save(pptx)
        pdf=root/'planeacion.pdf'; canvas=Canvas(str(pdf)); canvas.drawString(50,780,'Fecha | Actividad | TH a cargo'); canvas.drawString(50,760,'20/08/2026 | Encuentro familiar | Psicosocial'); canvas.save()
        previews={p.suffix:construir_preview_cronograma(str(p),p.name) for p in (xlsx,docx,pptx,pdf)}
        for ext,preview in previews.items():
            req(preview['actividades'],f'{ext} no detectó actividades')
            req('confianza' in preview['actividades'][0] and 'origen' in preview['actividades'][0],f'{ext} no informa confianza/origen')
        req(previews['.pdf']['requiere_revision'] is True,'PDF no exige revisión humana')
        repo=CalendarioInteligenteRepository(str(root/'db.sqlite3'),str(root/'uploads')); repo.init_schema(force=True)
        with tenant_context(1,role='COORDINADOR',username='coord'):
            preview=repo.registrar_preview_cronograma(str(xlsx),xlsx.name,'coord')
            with repo.connect() as conn: state=conn.execute('SELECT estado FROM calendario_cronogramas WHERE id=?',(preview['cronograma_id'],)).fetchone()['estado']
            req(state=='LISTO_PARA_REVISION','Importación no queda pendiente de aprobación humana')
            result=repo.confirmar_cronograma(preview['cronograma_id'],preview['actividades'],'coord')
            req(result['creados']==1,'Confirmación humana no creó la actividad')
            with repo.connect() as conn: state=conn.execute('SELECT estado FROM calendario_cronogramas WHERE id=?',(preview['cronograma_id'],)).fetchone()['estado']
            req(state=='APROBADO','Importación confirmada no quedó aprobada')
            idp_preview=repo.registrar_preview_actividades([{'fecha':'2026-08-28','actividad':'Entregar informe mensual','responsable':'Coordinación','entregable':'Informe y evidencias'}],'entregables agosto.jpeg','coord','IDP:25')
            req(idp_preview['actividades'][0]['fecha_limite']=='2026-08-28','La conexión IDP no conservó el día de entrega')
            idp_result=repo.confirmar_cronograma(idp_preview['cronograma_id'],idp_preview['actividades'],'coord')
            req(idp_result['creados']==1,'La actividad revisada del Motor Universal no llegó al calendario')
            req(any(item['fecha_limite']=='2026-08-28' and item['titulo']=='Entregar informe mensual' for item in repo.list_entregables({'periodo':'2026-08'})),'El día 28 no muestra la actividad importada desde IDP')
    print('PASS test_calendar_phase5_document_reader_v2_7_0')
if __name__=='__main__': run()
