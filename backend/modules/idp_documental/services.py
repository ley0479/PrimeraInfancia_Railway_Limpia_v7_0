from __future__ import annotations

from datetime import datetime
import hashlib
import json
import os
from pathlib import Path
import re
import unicodedata
import zipfile
import time
from difflib import SequenceMatcher
from typing import Any

from modules.dbapi_compat import sqlite3

from .schema import IDP_SCHEMA_SQL


ALLOWED_EXTENSIONS = {'.xlsx', '.xlsm', '.docx', '.pptx', '.pdf', '.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.heif', '.heic'}
MAX_FILE_SIZE = 50 * 1024 * 1024
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.heif', '.heic'}
AZURE_API_VERSION = '2024-11-30'


def now_iso() -> str:
    return datetime.now().isoformat(timespec='seconds')


def connect(database_path: str):
    conn = sqlite3.connect(database_path)
    conn.row_factory = sqlite3.Row
    return conn


def init_schema(database_path: str) -> None:
    conn = connect(database_path)
    conn.executescript(IDP_SCHEMA_SQL)
    conn.commit()
    conn.close()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open('rb') as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b''):
            digest.update(block)
    return digest.hexdigest()


def validate_file_signature(path: Path) -> None:
    ext = path.suffix.lower()
    if ext in {'.xlsx', '.xlsm', '.docx', '.pptx'}:
        if not zipfile.is_zipfile(path):
            raise ValueError('El archivo Office no tiene una estructura válida.')
        return
    if ext == '.pdf':
        with path.open('rb') as stream:
            if stream.read(5) != b'%PDF-':
                raise ValueError('El archivo no contiene una firma PDF válida.')
        return
    if ext in IMAGE_EXTENSIONS:
        try:
            from PIL import Image
            with Image.open(path) as image:
                image.verify()
        except Exception as exc:
            raise ValueError('La imagen está dañada o no corresponde al formato indicado.') from exc


def normalize(value: Any) -> str:
    text = unicodedata.normalize('NFKD', str(value or '').strip().lower())
    text = ''.join(char for char in text if not unicodedata.combining(char))
    return ' '.join(re.sub(r'[^a-z0-9]+', ' ', text).split())


def col_letter_to_index(value: str) -> int:
    result=0
    for char in str(value or '').upper():
        if 'A'<=char<='Z': result=result*26+ord(char)-64
    return result


def classify_document(text: str, filename: str) -> tuple[str, float, str]:
    sample = normalize(f'{filename} {text[:15000]}')
    normalized_filename = normalize(filename)
    # Un tablero o imagen mensual de entregables puede mencionar RAM y
    # listados como productos, pero no por eso contiene filas de asistentes.
    # El propósito explícito del nombre debe prevalecer sobre esas menciones.
    if 'entregable' in normalized_filename:
        return ('CRONOGRAMA', 0.97, 'nombre de archivo identificado como cronograma de entregables')
    filename_markers = {
        'LISTADO_ASISTENCIA': ('listado asistencia', 'asistencia'),
        'RAM': ('formato ram', ' ram '),
        'RPP': (' rpp ',),
        'BIENESTARINA': ('bienestarina',),
        'CRONOGRAMA': ('cronograma', 'entregables'),
        'PLANEACION_PEDAGOGICA': ('planeacion pedagogica', 'planeacion'),
        'PESO_TALLA': ('peso talla', 'valoracion nutricional'),
        'ACTA': ('acta',),
        'INFORME': ('informe',),
    }
    rules = [
        ('LISTADO_ASISTENCIA', ('listado de asistencia', 'firma del asistente', 'asistio', 'participantes'), 0.94),
        ('RAM', ('formato ram', 'registro de asistencia mensual', 'f27 mt1 pp'), 0.95),
        ('RPP', ('rpp', 'racion preparada', 'minuta patron'), 0.92),
        ('BIENESTARINA', ('bienestarina', 'entrega de alimento', 'lote'), 0.92),
        ('CRONOGRAMA', ('cronograma', 'fecha limite', 'responsable', 'actividad'), 0.88),
        ('PLANEACION_PEDAGOGICA', ('planeacion pedagogica', 'intencionalidad pedagogica', 'experiencia pedagogica'), 0.90),
        ('PESO_TALLA', ('peso kg', 'talla cm', 'valoracion nutricional', 'perimetro braquial'), 0.93),
        ('ACTA', ('acta', 'orden del dia', 'compromisos'), 0.86),
        ('INFORME', ('informe', 'objetivo', 'conclusiones'), 0.78),
    ]
    best = ('NO_CLASIFICADO', 0.0, 'sin coincidencias suficientes')
    for kind, tokens, base in rules:
        hits = sum(1 for token in tokens if token in sample)
        if hits:
            confidence = min(0.99, base - 0.14 + hits * 0.07)
            if any(marker.strip() in normalized_filename for marker in filename_markers.get(kind, ())):
                confidence = min(0.99, confidence + 0.15)
            if confidence > best[1]:
                best = (kind, confidence, f'{hits} regla(s) semantica(s)')
    return best


def _read_excel(path: Path) -> dict:
    from openpyxl import load_workbook
    workbook = load_workbook(path, read_only=True, data_only=False, keep_vba=path.suffix.lower() == '.xlsm')
    sheets, fragments = [], []
    for sheet in workbook.worksheets:
        rows = []
        for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
            values = [value.isoformat() if hasattr(value, 'isoformat') else value for value in row]
            if any(value not in (None, '') for value in values):
                rows.append({'fila': row_number, 'valores': values})
                fragments.extend(str(value) for value in values if value not in (None, ''))
            if len(rows) >= 5000:
                break
        sheets.append({'nombre': sheet.title, 'filas': rows, 'max_filas_leidas': len(rows)})
    workbook.close()
    return {'motor': 'OPENPYXL_NATIVE', 'texto': '\n'.join(fragments), 'hojas': sheets}


def _read_word(path: Path) -> dict:
    from docx import Document
    document = Document(str(path))
    paragraphs = [{'indice': index + 1, 'texto': p.text} for index, p in enumerate(document.paragraphs) if p.text.strip()]
    tables = []
    for table_index, table in enumerate(document.tables, 1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append({'tabla': table_index, 'filas': rows})
    text = '\n'.join([item['texto'] for item in paragraphs] + [str(cell) for table in tables for row in table['filas'] for cell in row])
    sheets=[{'nombre':f'TABLA_WORD_{table["tabla"]}','filas':[{'fila':index+1,'valores':row} for index,row in enumerate(table['filas'])],'max_filas_leidas':len(table['filas'])} for table in tables]
    return {'motor': 'PYTHON_DOCX_NATIVE', 'texto': text, 'parrafos': paragraphs, 'tablas': tables, 'hojas': sheets}


def _read_powerpoint(path: Path) -> dict:
    from pptx import Presentation
    presentation = Presentation(str(path))
    slides, fragments = [], []
    for slide_index, slide in enumerate(presentation.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and str(shape.text).strip():
                texts.append(str(shape.text))
                fragments.append(str(shape.text))
        slides.append({'pagina': slide_index, 'textos': texts})
    return {'motor': 'PYTHON_PPTX_NATIVE', 'texto': '\n'.join(fragments), 'diapositivas': slides}


def _read_pdf(path: Path) -> dict:
    try:
        import fitz
    except Exception as exc:
        return {'motor': 'PDF_PENDIENTE', 'texto': '', 'paginas': [], 'requiere_ocr': True, 'advertencia': f'PyMuPDF no disponible: {exc}'}
    document = fitz.open(str(path))
    pages, fragments = [], []
    for index, page in enumerate(document, 1):
        text = page.get_text('text') or ''
        pages.append({'pagina': index, 'texto': text})
        fragments.append(text)
    document.close()
    combined = '\n'.join(fragments).strip()
    return {'motor': 'PYMUPDF_NATIVE', 'texto': combined, 'paginas': pages, 'requiere_ocr': len(combined) < 20}


def _ocr_image_text(image) -> str:
    try:
        import pytesseract
        try:
            return pytesseract.image_to_string(image,lang='spa+eng')
        except Exception:
            return pytesseract.image_to_string(image)
    except Exception as exc:
        raise RuntimeError('OCR local no disponible. Instala y configura el ejecutable Tesseract.') from exc


def analyze_image_quality(image) -> dict:
    from PIL import ImageFilter,ImageStat
    sample=image.convert('L').copy(); sample.thumbnail((1200,1200))
    stats=ImageStat.Stat(sample); brightness=float(stats.mean[0]); contrast=float(stats.stddev[0])
    sharpness=float(ImageStat.Stat(sample.filter(ImageFilter.FIND_EDGES)).var[0])
    issues=[]
    if min(image.size)<600: issues.append('RESOLUCION_INSUFICIENTE')
    if contrast<4: issues.append('CONTRASTE_INSUFICIENTE')
    if brightness<10: issues.append('IMAGEN_DEMASIADO_OSCURA')
    if brightness>252: issues.append('IMAGEN_SOBREEXPUESTA')
    if sharpness<2: issues.append('POSIBLE_DESENFOQUE')
    critical=any(issue in issues for issue in {'RESOLUCION_INSUFICIENTE','CONTRASTE_INSUFICIENTE','IMAGEN_DEMASIADO_OSCURA','IMAGEN_SOBREEXPUESTA'})
    return {'ancho_px':image.width,'alto_px':image.height,'legible':not critical,'requiere_revision':bool(issues) or min(image.size)<1200,'brillo':round(brightness,2),'contraste':round(contrast,2),'nitidez_bordes':round(sharpness,2),'problemas':issues,'rechazo_automatico':critical}


def read_document_ocr(path: Path) -> dict:
    ext=path.suffix.lower(); pages=[]; fragments=[]
    if ext in IMAGE_EXTENSIONS:
        from PIL import Image,ImageOps
        with Image.open(path) as source:
            image=ImageOps.exif_transpose(source).convert('RGB')
            quality=analyze_image_quality(image)
            if quality['rechazo_automatico']:
                return {'motor':'CONTROL_CALIDAD','texto':'','requiere_ocr':True,'calidad':quality,'advertencia':'La fotografía no cumple la calidad mínima. Toma una imagen completa, enfocada, bien iluminada y con mayor resolución.'}
            prepared=ImageOps.autocontrast(image.convert('L')).convert('RGB')
            quality['correcciones_aplicadas']=['ORIENTACION_EXIF','AUTOCONTRASTE']
            text=_ocr_image_text(prepared); fragments.append(text); pages.append({'pagina':1,'texto':text})
    elif ext=='.pdf':
        try: import fitz
        except Exception as exc: raise RuntimeError('PyMuPDF no está disponible para convertir el PDF escaneado.') from exc
        document=fitz.open(str(path))
        try:
            for index,page in enumerate(document,1):
                pixmap=page.get_pixmap(matrix=fitz.Matrix(2,2),alpha=False)
                from PIL import Image
                image=Image.frombytes('RGB',(pixmap.width,pixmap.height),pixmap.samples)
                text=_ocr_image_text(image); fragments.append(text); pages.append({'pagina':index,'texto':text})
        finally: document.close()
        quality={'paginas':len(pages),'requiere_revision':True}
    else:
        raise ValueError('El reintento OCR solo admite imágenes y PDF.')
    combined='\n'.join(fragments).strip()
    if len(combined)<20:
        return {'motor':'TESSERACT_LOCAL','texto':combined,'paginas':pages,'requiere_ocr':True,'calidad':quality,'advertencia':'El motor OCR no encontró texto suficiente. Revisa enfoque, iluminación y recorte.'}
    return {'motor':'TESSERACT_LOCAL','texto':combined,'paginas':pages,'requiere_ocr':False,'calidad':quality,'origen_ocr':True}


def azure_document_intelligence_configured() -> bool:
    return bool(str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT') or '').strip() and str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_KEY') or '').strip())


def _azure_result_to_raw(payload: dict) -> dict:
    result=payload.get('analyzeResult') or {}; content=str(result.get('content') or '')
    pages=[]
    for index,page in enumerate(result.get('pages') or [],1):
        lines=page.get('lines') or []; pages.append({'pagina':page.get('pageNumber') or index,'texto':'\n'.join(str(line.get('content') or '') for line in lines),'ancho':page.get('width'),'alto':page.get('height'),'unidad':page.get('unit')})
    sheets=[]
    for table_index,table in enumerate(result.get('tables') or [],1):
        row_count=int(table.get('rowCount') or 0); column_count=int(table.get('columnCount') or 0); matrix=[[None for _ in range(column_count)] for _ in range(row_count)]; confidences=[[None for _ in range(column_count)] for _ in range(row_count)]
        for cell in table.get('cells') or []:
            row=int(cell.get('rowIndex') or 0); column=int(cell.get('columnIndex') or 0)
            if row<row_count and column<column_count: matrix[row][column]=cell.get('content'); confidences[row][column]=cell.get('confidence')
        sheets.append({'nombre':f'TABLA_OCR_{table_index}','filas':[{'fila':row+1,'valores':values,'confianzas':confidences[row]} for row,values in enumerate(matrix)],'max_filas_leidas':row_count})
    return {'motor':'AZURE_DOCUMENT_INTELLIGENCE','texto':content or '\n'.join(page['texto'] for page in pages),'paginas':pages,'hojas':sheets,'requiere_ocr':False,'origen_ocr':True,'modelo_azure':result.get('modelId'),'api_version':result.get('apiVersion') or AZURE_API_VERSION}


def read_document_azure(path: Path) -> dict:
    endpoint=str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT') or '').strip().rstrip('/'); key=str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_KEY') or '').strip()
    if not endpoint or not key: raise RuntimeError('Azure Document Intelligence no está configurado.')
    import requests
    model=str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_MODEL') or 'prebuilt-layout').strip(); api_version=str(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_API_VERSION') or AZURE_API_VERSION).strip(); timeout=max(15,int(os.environ.get('AZURE_DOCUMENT_INTELLIGENCE_TIMEOUT_SECONDS','120')))
    url=f'{endpoint}/documentintelligence/documentModels/{model}:analyze?api-version={api_version}'
    mime={'.pdf':'application/pdf','.png':'image/png','.jpg':'image/jpeg','.jpeg':'image/jpeg','.bmp':'image/bmp','.tif':'image/tiff','.tiff':'image/tiff','.heif':'image/heif','.heic':'image/heif'}.get(path.suffix.lower(),'application/octet-stream')
    with path.open('rb') as stream: response=requests.post(url,headers={'Ocp-Apim-Subscription-Key':key,'Content-Type':mime},data=stream,timeout=30)
    if response.status_code!=202: raise RuntimeError(f'Azure rechazó el documento (HTTP {response.status_code}).')
    operation=response.headers.get('Operation-Location')
    if not operation: raise RuntimeError('Azure no devolvió la ubicación del resultado.')
    deadline=time.monotonic()+timeout
    while time.monotonic()<deadline:
        result_response=requests.get(operation,headers={'Ocp-Apim-Subscription-Key':key},timeout=30)
        if result_response.status_code>=400: raise RuntimeError(f'Azure no pudo consultar el resultado (HTTP {result_response.status_code}).')
        payload=result_response.json(); status=str(payload.get('status') or '').lower()
        if status=='succeeded': return _azure_result_to_raw(payload)
        if status in {'failed','canceled'}: raise RuntimeError('Azure no pudo analizar el documento.')
        time.sleep(1)
    raise TimeoutError('Azure Document Intelligence excedió el tiempo máximo de análisis.')


def read_document_intelligent(path: Path) -> dict:
    native=read_document(path)
    if not native.get('requiere_ocr'): return native
    if azure_document_intelligence_configured():
        try: return read_document_azure(path)
        except Exception as exc:
            fallback=read_document_ocr(path); fallback['advertencia_azure']=str(exc)[:300]; return fallback
    if path.suffix.lower() in IMAGE_EXTENSIONS or path.suffix.lower()=='.pdf': return read_document_ocr(path)
    return native


def read_document(path: Path) -> dict:
    ext = path.suffix.lower()
    if ext in {'.xlsx', '.xlsm'}:
        return _read_excel(path)
    if ext == '.docx':
        return _read_word(path)
    if ext == '.pptx':
        return _read_powerpoint(path)
    if ext == '.pdf':
        return _read_pdf(path)
    if ext in IMAGE_EXTENSIONS:
        quality = {'legible': None, 'requiere_revision': True}
        try:
            from PIL import Image,ImageOps
            with Image.open(path) as image:
                corrected=ImageOps.exif_transpose(image)
                quality=analyze_image_quality(corrected); quality['modo']=corrected.mode
        except Exception:
            pass
        return {'motor': 'IMAGEN_PENDIENTE_OCR', 'texto': '', 'requiere_ocr': True, 'calidad': quality}
    raise ValueError('Formato no soportado por el lector IDP.')


HEADER_ALIASES = {
    'documento': ('documento', 'identificacion', 'cedula', 'nui'),
    'nombre_completo': ('nombre', 'nombres y apellidos', 'participante', 'beneficiario'),
    'asistio': ('asistio', 'asistencia', 'presente'),
    'firma_presente': ('firma', 'firma del asistente'),
    'unidad': ('uds', 'uca', 'unidad', 'unidad de servicio'),
}
RAM_HEADER_ALIASES = {
    'tipo_documento':('tipo documento','tipo de documento','tipo doc'),
    'documento':('documento','documento beneficiario','identificacion','cedula','nui'),
    'primer_nombre':('primer nombre',),
    'segundo_nombre':('segundo nombre',),
    'primer_apellido':('primer apellido',),
    'segundo_apellido':('segundo apellido',),
    'nombre_completo':('nombre','nombre completo','nombres y apellidos','beneficiario'),
    'unidad':('uds','uca','unidad','unidad de servicio'),
    'total_asistencias':('total asistencias','total asistencia'),
    'total_inasistencias':('total inasistencias','inasistencias'),
    'causa_retiro':('causa retiro','causa de retiro'),
}
BIENESTARINA_HEADER_ALIASES = {
    'tipo_documento':('tipo documento','tipo de documento','tipo doc'),
    'documento':('documento','nui','identificacion','documento beneficiario'),
    'primer_nombre':('primer nombre',), 'segundo_nombre':('segundo nombre',),
    'primer_apellido':('primer apellido',), 'segundo_apellido':('segundo apellido',),
    'nombre_completo':('nombre','nombre completo','beneficiario'),
    'fecha_entrega':('fecha entrega','fecha de entrega'),
    'lote':('lote','numero de lote'),
    'cantidad':('cantidad','cantidad entregada','cantidad bienestarina'),
    'acudiente':('acudiente','nombre acudiente'),
    'parentesco':('parentesco',), 'firma_presente':('firma','firma acudiente','firma de recibido'),
    'unidad':('uds','uca','unidad','unidad de servicio'),
}
SCHEDULE_HEADER_ALIASES = {
    'fecha':('fecha','fecha actividad','fecha de actividad','fecha programada'),
    'fecha_limite':('fecha limite','fecha de entrega','entrega','vence','vencimiento'),
    'actividad':('actividad','tema','tarea','descripcion','compromiso'),
    'responsable':('responsable','th a cargo','profesional','encargado'),
    'entregable':('entregable','evidencia','producto','soporte'),
    'modulo':('modulo','componente','area'),
    'periodicidad':('periodicidad','frecuencia','recurrencia'),
}
NUTRITION_HEADER_ALIASES = {
    'documento':('documento','identificacion','cedula','nui'),
    'nombre_completo':('nombre','nombre completo','nombres y apellidos','beneficiario','participante'),
    'fecha':('fecha','fecha valoracion','fecha de valoracion'),
    'peso_kg':('peso','peso kg','peso en kg','peso kilogramo'),
    'talla_cm':('talla','talla cm','estatura','longitud cm'),
    'perimetro_braquial_cm':('perimetro braquial','perimetro braquial cm','pb cm'),
    'unidad':('uds','uca','unidad','unidad de servicio'),
}
PLANNING_LABEL_ALIASES = {
    'unidad':('uds','uca','unidad','unidad de servicio'),
    'periodo':('periodo','mes','vigencia'),
    'tema':('tema','nombre de la experiencia','titulo'),
    'objetivo':('objetivo','intencionalidad','intencionalidad pedagogica','proposito'),
    'actividad':('actividad','experiencia pedagogica','experiencia','desarrollo'),
    'fecha_programada':('fecha','fecha programada','fecha de realizacion'),
    'poblacion_objetivo':('poblacion objetivo','grupo etario','participantes'),
    'evidencia_requerida':('evidencia','evidencia requerida','entregables','soportes'),
    'tipo_encuentro':('tipo de encuentro','tipo actividad','modalidad'),
    'responsable':('responsable','docente','agente educativo','th a cargo'),
    'recursos':('recursos','materiales'),
    'observaciones':('observaciones','valoracion','evaluacion'),
}
MINUTES_LABEL_ALIASES = {
    'fecha':('fecha','fecha del acta'), 'lugar':('lugar','sitio'), 'hora':('hora','hora de inicio'),
    'tema':('tema','asunto','objetivo'), 'asistentes':('asistentes','participantes'),
    'orden_dia':('orden del dia','agenda'), 'desarrollo':('desarrollo','desarrollo de la reunion'),
    'compromisos':('compromisos','acuerdos'), 'responsable':('responsable','elaboro','convoca'),
    'firma_referencia':('firmas','firma','firma de asistentes'), 'unidad':('uds','uca','unidad','unidad de servicio'),
}
REPORT_LABEL_ALIASES = {
    'periodo':('periodo','mes','vigencia'), 'fecha':('fecha','fecha del informe'), 'titulo':('titulo','nombre del informe','tema'),
    'objetivo':('objetivo','objetivo general'), 'actividades':('actividades','actividades realizadas'),
    'resultados':('resultados','logros','avance'), 'dificultades':('dificultades','situaciones encontradas'),
    'conclusiones':('conclusiones','conclusion'), 'recomendaciones':('recomendaciones','acciones de mejora'),
    'responsable':('responsable','elaborado por','profesional'), 'unidad':('uds','uca','unidad','unidad de servicio'),
}
RPP_LABEL_ALIASES = {
    'fecha':('fecha','fecha de preparacion','fecha del servicio'),
    'periodo':('periodo','mes','vigencia'),
    'unidad':('uds','uca','unidad','unidad de servicio'),
    'modalidad':('modalidad','modalidad de atencion'),
    'tiempo_comida':('tiempo de comida','servicio de alimentacion','momento de consumo'),
    'preparacion':('preparacion','menu','nombre de la preparacion'),
    'minuta_patron':('minuta patron','ciclo de minuta','minuta'),
    'porciones':('porciones','numero de porciones','raciones preparadas','cantidad de raciones'),
    'responsable':('responsable','manipulador de alimentos','preparado por'),
    'observaciones':('observaciones','novedades'),
}


def _mapped_header(value: Any) -> str | None:
    text = normalize(value)
    for field, aliases in HEADER_ALIASES.items():
        if text in aliases or any(len(alias) >= 4 and alias in text for alias in aliases):
            return field
    return None


def _mapped_ram_header(value: Any) -> str | None:
    text=normalize(value)
    for field,aliases in RAM_HEADER_ALIASES.items():
        if text in aliases: return field
    return None


def _mapped_bienestarina_header(value: Any) -> str | None:
    text=normalize(value)
    for field,aliases in BIENESTARINA_HEADER_ALIASES.items():
        if text in aliases: return field
    return None


def _ram_day_header(value: Any) -> int | None:
    text=normalize(value)
    match=re.fullmatch(r'(?:dia )?([1-9]|[12][0-9]|3[01])',text)
    return int(match.group(1)) if match else None


def _mapped_schedule_header(value: Any) -> str | None:
    text=normalize(value)
    for field,aliases in SCHEDULE_HEADER_ALIASES.items():
        if text in aliases or any(len(alias)>=4 and alias in text for alias in aliases): return field
    return None


def _mapped_nutrition_header(value: Any) -> str | None:
    text=normalize(value)
    for field,aliases in NUTRITION_HEADER_ALIASES.items():
        if text in aliases or any(len(alias)>=4 and alias in text for alias in aliases): return field
    return None


def _numeric_measure(value: Any) -> float | None:
    if value in (None,''): return None
    match=re.search(r'-?\d+(?:[.,]\d+)?',str(value).replace(' ','').replace(',','.'))
    try: return float(match.group(0)) if match else None
    except ValueError: return None


def _planning_label(value: Any) -> str | None:
    text=normalize(value).rstrip(':')
    for field,aliases in PLANNING_LABEL_ALIASES.items():
        if text in aliases or any(len(alias)>=5 and text.startswith(alias) for alias in aliases): return field
    return None


def _label_for(value: Any, aliases: dict[str,tuple[str,...]]) -> str | None:
    text=normalize(value).rstrip(':')
    for field,names in aliases.items():
        if text in names or any(len(name)>=5 and text.startswith(name) for name in names): return field
    return None


def _schedule_date(value: Any) -> str | None:
    if value in (None,''): return None
    if hasattr(value,'strftime'): return value.strftime('%Y-%m-%d')
    text=str(value).strip()
    for pattern in ('%Y-%m-%d','%d/%m/%Y','%d-%m-%Y','%Y/%m/%d'):
        try: return datetime.strptime(text[:10],pattern).date().isoformat()
        except ValueError: pass
    return None


def _canonicalize_schedule(raw: dict, canonical: dict, fields: list[dict]) -> None:
    canonical['actividades']=[]
    for sheet in raw.get('hojas') or []:
        rows=sheet.get('filas') or []
        for position,row in enumerate(rows):
            mapping={field:column for column,value in enumerate(row.get('valores') or []) if (field:=_mapped_schedule_header(value))}
            if 'actividad' not in mapping or not ({'fecha','fecha_limite','responsable','entregable'} & set(mapping)): continue
            for data_row in rows[position+1:]:
                values=data_row.get('valores') or []; activity={}; index=len(canonical['actividades'])
                for field,column in mapping.items():
                    raw_value=values[column] if column<len(values) else None; value=_schedule_date(raw_value) if field in {'fecha','fecha_limite'} else (str(raw_value).strip() if raw_value not in (None,'') else None)
                    if value is not None: activity[field]=value
                    fields.append({'ruta':f'actividades.{index}.{field}','valor':value,'texto_original':raw_value,'confianza':.96 if value is not None and not raw.get('origen_ocr') else (.76 if value is not None else 0),'evidencia':{'hoja':sheet['nombre'],'fila':data_row.get('fila'),'columna':column+1},'regla':'encabezado_cronograma'})
                if activity.get('actividad'): canonical['actividades'].append(activity)
            return


def _canonicalize_nutrition(raw: dict, canonical: dict, fields: list[dict]) -> None:
    canonical['valoraciones']=[]
    for sheet in raw.get('hojas') or []:
        rows=sheet.get('filas') or []
        for position,row in enumerate(rows):
            mapping={field:column for column,value in enumerate(row.get('valores') or []) if (field:=_mapped_nutrition_header(value))}
            if not ({'documento','nombre_completo'} & set(mapping)) or not ({'peso_kg','talla_cm'} & set(mapping)): continue
            for data_row in rows[position+1:]:
                values=data_row.get('valores') or []; record={}; index=len(canonical['valoraciones'])
                for field,column in mapping.items():
                    raw_value=values[column] if column<len(values) else None
                    if field in {'peso_kg','talla_cm','perimetro_braquial_cm'}: value=_numeric_measure(raw_value)
                    elif field=='fecha': value=_schedule_date(raw_value)
                    else: value=str(raw_value).strip() if raw_value not in (None,'') else None
                    if value is not None: record[field]=value
                    fields.append({'ruta':f'valoraciones.{index}.{field}','valor':value,'texto_original':raw_value,'confianza':.96 if value is not None and not raw.get('origen_ocr') else (.76 if value is not None else 0),'evidencia':{'hoja':sheet['nombre'],'fila':data_row.get('fila'),'columna':column+1},'regla':'encabezado_peso_talla'})
                if record.get('documento') or record.get('nombre_completo'):
                    canonical['valoraciones'].append(record)
                    if record.get('unidad'): canonical['unidad_servicio'].setdefault('nombre',record['unidad'])
            return


def _canonicalize_ram(raw: dict, canonical: dict, fields: list[dict]) -> None:
    canonical['participantes']=[]
    for sheet in raw.get('hojas') or []:
        rows=sheet.get('filas') or []
        for position,row in enumerate(rows):
            values=row.get('valores') or []
            mapping={field:column for column,value in enumerate(values) if (field:=_mapped_ram_header(value))}
            day_columns={day:column for column,value in enumerate(values) if (day:=_ram_day_header(value))}
            if 'documento' not in mapping or not ({'nombre_completo','primer_nombre','primer_apellido'} & set(mapping)): continue
            for data_row in rows[position+1:]:
                row_values=data_row.get('valores') or []; participant={}; index=len(canonical['participantes'])
                for field,column in mapping.items():
                    value=row_values[column] if column<len(row_values) else None
                    if value in (None,''): continue
                    if field=='unidad': canonical['unidad_servicio'].setdefault('nombre',str(value).strip()); continue
                    participant[field]=str(value).strip()
                    fields.append({'ruta':f'participantes.{index}.{field}','valor':participant[field],'texto_original':value,'confianza':.98,'evidencia':{'hoja':sheet['nombre'],'fila':data_row.get('fila'),'columna':column+1},'regla':'encabezado_ram'})
                name_parts=[participant.pop(field,None) for field in ('primer_nombre','segundo_nombre','primer_apellido','segundo_apellido')]
                if not participant.get('nombre_completo') and any(name_parts): participant['nombre_completo']=' '.join(part for part in name_parts if part)
                attendance={}
                for day,column in day_columns.items():
                    value=row_values[column] if column<len(row_values) else None
                    if value not in (None,''): attendance[str(day)]=str(value).strip().upper()
                if day_columns:
                    participant['asistencia_dias']=attendance
                    fields.append({'ruta':f'participantes.{index}.asistencia_dias','valor':attendance,'texto_original':attendance,'confianza':.98,'evidencia':{'hoja':sheet['nombre'],'fila':data_row.get('fila'),'columnas':len(day_columns)},'regla':'dias_ram'})
                if participant.get('documento') or participant.get('nombre_completo'): canonical['participantes'].append(participant)
            return


def _canonicalize_bienestarina(raw: dict, canonical: dict, fields: list[dict]) -> None:
    canonical['entregas']=[]
    for sheet in raw.get('hojas') or []:
        rows=sheet.get('filas') or []
        for position,row in enumerate(rows):
            mapping={field:column for column,value in enumerate(row.get('valores') or []) if (field:=_mapped_bienestarina_header(value))}
            if 'documento' not in mapping or not ({'nombre_completo','primer_nombre','primer_apellido'} & set(mapping)): continue
            for data_row in rows[position+1:]:
                values=data_row.get('valores') or []; delivery={}; index=len(canonical['entregas'])
                for field,column in mapping.items():
                    raw_value=values[column] if column<len(values) else None
                    if raw_value in (None,''): continue
                    if field=='unidad': canonical['unidad_servicio'].setdefault('nombre',str(raw_value).strip()); continue
                    value=_schedule_date(raw_value) if field=='fecha_entrega' else str(raw_value).strip()
                    delivery[field]=value
                    fields.append({'ruta':f'entregas.{index}.{field}','valor':value,'texto_original':raw_value,'confianza':.98,'evidencia':{'hoja':sheet['nombre'],'fila':data_row.get('fila'),'columna':column+1},'regla':'encabezado_bienestarina'})
                name_parts=[delivery.pop(field,None) for field in ('primer_nombre','segundo_nombre','primer_apellido','segundo_apellido')]
                if not delivery.get('nombre_completo') and any(name_parts): delivery['nombre_completo']=' '.join(part for part in name_parts if part)
                if delivery.get('documento') or delivery.get('nombre_completo'):
                    canonical['entregas'].append(delivery)
                    canonical['participantes'].append({field:delivery[field] for field in ('documento','nombre_completo') if delivery.get(field)})
            return


def _canonicalize_planning(raw: dict, canonical: dict, fields: list[dict]) -> None:
    planning={}; evidence={}
    paragraphs=raw.get('parrafos') or []
    for position,item in enumerate(paragraphs):
        text=str(item.get('texto') or '').strip(); label_part,value_part=(text.split(':',1)+[''])[:2] if ':' in text else (text,'')
        field=_planning_label(label_part)
        if not field: continue
        value=value_part.strip()
        if not value and position+1<len(paragraphs) and not _planning_label(paragraphs[position+1].get('texto')): value=str(paragraphs[position+1].get('texto') or '').strip()
        if value: planning[field]=_schedule_date(value) if field=='fecha_programada' else value; evidence[field]={'parrafo':item.get('indice'),'texto':text}
    for sheet in raw.get('hojas') or []:
        rows=sheet.get('filas') or []
        for row in rows:
            values=row.get('valores') or []
            if len(values)>=2:
                field=_planning_label(values[0]); value=values[1]
                if field and value not in (None,'') and field not in planning: planning[field]=_schedule_date(value) if field=='fecha_programada' else str(value).strip(); evidence[field]={'hoja':sheet['nombre'],'fila':row.get('fila'),'columna':2}
        for position,row in enumerate(rows):
            mapping={field:column for column,value in enumerate(row.get('valores') or []) if (field:=_planning_label(value))}
            if len(mapping)<2 or position+1>=len(rows): continue
            values=rows[position+1].get('valores') or []
            for field,column in mapping.items():
                value=values[column] if column<len(values) else None
                if value not in (None,'') and field not in planning: planning[field]=_schedule_date(value) if field=='fecha_programada' else str(value).strip(); evidence[field]={'hoja':sheet['nombre'],'fila':rows[position+1].get('fila'),'columna':column+1}
            break
    canonical['planeacion']=planning
    if planning.get('unidad'): canonical['unidad_servicio']['nombre']=planning['unidad']
    for field,value in planning.items(): fields.append({'ruta':f'planeacion.{field}','valor':value,'texto_original':value,'confianza':.94 if not raw.get('origen_ocr') else .74,'evidencia':evidence.get(field) or {},'regla':'etiqueta_planeacion'})


def _canonicalize_labeled_document(raw: dict, canonical: dict, fields: list[dict], key: str, aliases: dict[str,tuple[str,...]]) -> None:
    data={}; evidence={}; paragraphs=raw.get('parrafos') or []
    for position,item in enumerate(paragraphs):
        text=str(item.get('texto') or '').strip(); label_part,value_part=(text.split(':',1)+[''])[:2] if ':' in text else (text,''); field=_label_for(label_part,aliases)
        if not field: continue
        value=value_part.strip()
        if not value and position+1<len(paragraphs) and not _label_for(paragraphs[position+1].get('texto'),aliases): value=str(paragraphs[position+1].get('texto') or '').strip()
        if value: data[field]=_schedule_date(value) if field=='fecha' else value; evidence[field]={'parrafo':item.get('indice'),'texto':text}
    for sheet in raw.get('hojas') or []:
        for row in sheet.get('filas') or []:
            values=row.get('valores') or []
            if len(values)<2: continue
            field=_label_for(values[0],aliases); value=values[1]
            if field and value not in (None,'') and field not in data: data[field]=_schedule_date(value) if field=='fecha' else str(value).strip(); evidence[field]={'hoja':sheet['nombre'],'fila':row.get('fila'),'columna':2}
    canonical[key]=data
    if data.get('unidad'): canonical['unidad_servicio']['nombre']=data['unidad']
    for field,value in data.items(): fields.append({'ruta':f'{key}.{field}','valor':value,'texto_original':value,'confianza':.93 if not raw.get('origen_ocr') else .73,'evidencia':evidence.get(field) or {},'regla':f'etiqueta_{key}'})


def _canonicalize_ocr_attendance(raw: dict, canonical: dict, fields: list[dict]) -> None:
    seen=set()
    for line_number,line in enumerate(str(raw.get('texto') or '').splitlines(),1):
        original=' '.join(line.split())
        if not original or any(token in normalize(original) for token in ('listado de asistencia','nombre documento','nombre completo documento')):
            continue
        match=re.search(r'(?<!\d)(\d{4,12})(?!\d)',original)
        if not match: continue
        document=match.group(1)
        if document in seen or (len(document)==4 and 1900<=int(document)<=2100): continue
        prefix=re.sub(r'^\s*\d+[.)-]?\s*','',original[:match.start()]).strip(' :-|')
        suffix=original[match.end():].strip(' :-|')
        if len(normalize(prefix))<3: continue
        participant={'nombre_completo':prefix,'documento':document}
        unit_match=re.search(r'\b(?:UDS|UCA)\s*[:#-]?\s*[A-Za-z0-9._-]+',suffix,re.IGNORECASE)
        if unit_match:
            participant['unidad']=unit_match.group(0).strip()
            canonical['unidad_servicio'].setdefault('nombre',participant['unidad'])
        attendance_match=re.search(r'\b(SI|NO|PRESENTE|AUSENTE)\b',suffix,re.IGNORECASE)
        if attendance_match: participant['asistio']=normalize(attendance_match.group(1)) in {'si','presente'}
        index=len(canonical['participantes']); canonical['participantes'].append(participant); seen.add(document)
        for field,value in participant.items():
            confidence=.82 if field=='documento' else (.74 if field=='nombre_completo' else .68)
            fields.append({'ruta':f'participantes.{index}.{field}','valor':value,'texto_original':original,'confianza':confidence,'evidencia':{'pagina':1,'linea':line_number,'texto':original},'regla':'fila_ocr_con_documento'})


def canonicalize(raw: dict, document_type: str) -> tuple[dict, list[dict]]:
    canonical = {'tipo_documento': document_type, 'version_plantilla': None, 'periodo': {}, 'fundacion': {}, 'unidad_servicio': {}, 'actividad': {}, 'participantes': [], 'metadatos': {'motor': raw.get('motor'), 'requiere_ocr': bool(raw.get('requiere_ocr'))}}
    fields = []
    if document_type == 'LISTADO_ASISTENCIA' and raw.get('hojas'):
        for sheet in raw['hojas']:
            rows = sheet.get('filas') or []
            for position, row in enumerate(rows):
                mapping = {field: col for col, value in enumerate(row['valores']) if (field := _mapped_header(value))}
                if len(mapping) < 2:
                    continue
                for data_row in rows[position + 1:]:
                    values = data_row['valores']
                    participant = {}
                    for field, col in mapping.items():
                        value = values[col] if col < len(values) else None
                        row_confidences=data_row.get('confianzas') or []; field_confidence=float(row_confidences[col] or 0.82) if col<len(row_confidences) else (0.98 if value not in (None,'') else 0.0)
                        if field == 'unidad' and value and not canonical['unidad_servicio'].get('nombre'):
                            canonical['unidad_servicio']['nombre'] = str(value)
                        elif field not in {'unidad'}:
                            participant[field] = value
                            fields.append({'ruta': f'participantes.{len(canonical["participantes"])}.{field}', 'valor': value, 'texto_original': value, 'confianza': field_confidence, 'evidencia': {'hoja': sheet['nombre'], 'fila': data_row['fila'], 'columna': col + 1}, 'regla': 'tabla_azure' if raw.get('motor')=='AZURE_DOCUMENT_INTELLIGENCE' else 'encabezado_excel'})
                    if any(participant.get(key) not in (None, '') for key in ('documento', 'nombre_completo')):
                        canonical['participantes'].append(participant)
                break
    if document_type=='RAM': _canonicalize_ram(raw,canonical,fields)
    if document_type=='BIENESTARINA': _canonicalize_bienestarina(raw,canonical,fields)
    if document_type in {'LISTADO_ASISTENCIA','RAM'} and raw.get('origen_ocr') and not canonical['participantes']:
        _canonicalize_ocr_attendance(raw,canonical,fields)
    if document_type=='CRONOGRAMA': _canonicalize_schedule(raw,canonical,fields)
    if document_type=='PESO_TALLA': _canonicalize_nutrition(raw,canonical,fields)
    if document_type=='PLANEACION_PEDAGOGICA': _canonicalize_planning(raw,canonical,fields)
    if document_type=='ACTA': _canonicalize_labeled_document(raw,canonical,fields,'acta',MINUTES_LABEL_ALIASES)
    if document_type=='INFORME': _canonicalize_labeled_document(raw,canonical,fields,'informe',REPORT_LABEL_ALIASES)
    if document_type=='RPP': _canonicalize_labeled_document(raw,canonical,fields,'rpp',RPP_LABEL_ALIASES)
    fields.append({'ruta': 'tipo_documento', 'valor': document_type, 'texto_original': document_type, 'confianza': 1.0, 'evidencia': {}, 'regla': 'clasificador_reglas'})
    return canonical, fields


def _document_key(value: Any) -> str:
    return re.sub(r'[^A-Za-z0-9]', '', str(value or '')).upper()


def resolve_official_template_version(database_path: str, tenant_id: int, document_type: str) -> dict | None:
    aliases={
        'LISTADO_ASISTENCIA':('LISTADO_ASISTENCIA','LISTADO_ASISTENCIA_USUARIOS','ASISTENCIA_USUARIOS'),
        'RAM':('RAM','RAN','RRAN'),
        'RPP':('RPP','RPP','RPP'),
        'BIENESTARINA':('BIENESTARINA','BIENESTARINA','BIENESTARINA'),
    }.get(str(document_type or '').upper(),(str(document_type or '').upper(),)*3)
    try:
        conn=connect(database_path)
        row=conn.execute("""SELECT id,tipo_formato,codigo,nombre,version,fecha_vigencia,fecha_vigencia_fin,estado,hash_sha256,mapeo_json FROM plantillas_oficiales_versiones WHERE COALESCE(fundacion_id,1)=? AND UPPER(tipo_formato) IN (?,?,?) AND LOWER(COALESCE(estado,'')) IN ('vigente','activa','publicada') ORDER BY COALESCE(fecha_vigencia,'') DESC,COALESCE(updated_at,'') DESC,id DESC LIMIT 1""",(tenant_id,*aliases)).fetchone(); conn.close()
    except Exception:
        return None
    if not row: return None
    item=dict(row)
    try: mapping=json.loads(item.pop('mapeo_json',None) or '{}')
    except Exception: mapping={}
    if isinstance(mapping,dict) and isinstance(mapping.get('campos'),dict):
        mapped_items=[{'field':field,'col':column,'sheet':mapping.get('hoja'),'data_start_row':mapping.get('fila_datos')} for field,column in mapping['campos'].items()]
    else: mapped_items=(mapping.get('fields') or mapping.get('campos') or mapping) if isinstance(mapping,dict) else mapping
    item['mapeo_resumen']={'campos':len(mapped_items) if isinstance(mapped_items,(list,dict)) else 0}
    item['_mapeo']=mapped_items
    return item


def apply_official_mapping(raw: dict, canonical: dict, fields: list[dict], template_version: dict | None) -> None:
    """Usa el mapeo publicado solo como fallback; nunca sobrescribe extracción existente."""
    kind=canonical.get('tipo_documento')
    if kind not in {'RAM','LISTADO_ASISTENCIA','BIENESTARINA'} or not template_version: return
    if kind in {'RAM','LISTADO_ASISTENCIA'} and canonical.get('participantes'): return
    if kind=='BIENESTARINA' and canonical.get('entregas'): return
    mapping=template_version.get('_mapeo')
    if not isinstance(mapping,list) or not mapping: return
    sheets={normalize(sheet.get('nombre')):sheet for sheet in raw.get('hojas') or []}
    target_name=normalize(next((item.get('sheet') or item.get('hoja') for item in mapping if item.get('sheet') or item.get('hoja')),''))
    sheet=sheets.get(target_name) or next(iter(sheets.values()),None)
    if not sheet: return
    rows={int(row.get('fila') or 0):row.get('valores') or [] for row in sheet.get('filas') or []}
    start=min((int(item.get('data_start_row') or item.get('fila_inicio') or 0) for item in mapping if item.get('data_start_row') or item.get('fila_inicio')),default=0)
    configured_ends=[int(item.get('fila_fin')) for item in mapping if item.get('fila_fin')]
    end=max(configured_ends,default=max(rows,default=start))
    aliases={'documento_beneficiario':'documento','documento':'documento','nui':'documento','tipo_documento':'tipo_documento','primer_nombre':'primer_nombre','segundo_nombre':'segundo_nombre','primer_apellido':'primer_apellido','segundo_apellido':'segundo_apellido','nombre_completo':'nombre_completo','nombre':'nombre_completo','total_asistencias':'total_asistencias','total_inasistencias':'total_inasistencias','causa_retiro':'causa_retiro','asistencia':'asistio','asistio':'asistio','firma':'firma_presente','firma_presente':'firma_presente','unidad':'unidad','unidad_servicio':'unidad','fecha_entrega':'fecha_entrega','lote':'lote','cantidad':'cantidad','nombre_acudiente':'acudiente','acudiente':'acudiente','parentesco':'parentesco'}
    for row_number in range(start,end+1):
        values=rows.get(row_number) or []; record={}; index=len(canonical.get('entregas') or []) if kind=='BIENESTARINA' else len(canonical['participantes'])
        for item in mapping:
            source=str(item.get('field') or item.get('campo') or ''); field=aliases.get(source)
            column=int(item.get('col') or item.get('col_index') or 0)
            if source=='control_asistencia':
                span=str(item.get('col_letter') or item.get('columna') or '')
                bounds=span.split(':'); first=column; last=col_letter_to_index(bounds[-1]) if len(bounds)>1 else column
                attendance={str(day):str(values[col-1]).strip().upper() for day,col in enumerate(range(first,last+1),1) if col-1<len(values) and values[col-1] not in (None,'')}
                if attendance: record['asistencia_dias']=attendance
                continue
            if not field or column<1 or column-1>=len(values) or values[column-1] in (None,''): continue
            value=_schedule_date(values[column-1]) if field=='fecha_entrega' else str(values[column-1]).strip()
            if field in {'asistio','firma_presente'}: value=normalize(value) in {'si','x','presente','asistio','firmado'}
            if field=='unidad': canonical['unidad_servicio'].setdefault('nombre',str(value)); continue
            record[field]=value
            route_root='entregas' if kind=='BIENESTARINA' else 'participantes'
            fields.append({'ruta':f'{route_root}.{index}.{field}','valor':value,'texto_original':values[column-1],'confianza':.99,'evidencia':{'hoja':sheet.get('nombre'),'fila':row_number,'columna':column,'plantilla_version_id':template_version.get('id')},'regla':'mapeo_oficial_versionado'})
        name_parts=[record.pop(field,None) for field in ('primer_nombre','segundo_nombre','primer_apellido','segundo_apellido')]
        if not record.get('nombre_completo') and any(name_parts): record['nombre_completo']=' '.join(part for part in name_parts if part)
        if not (record.get('documento') or record.get('nombre_completo')): continue
        if kind=='BIENESTARINA':
            canonical.setdefault('entregas',[]).append(record)
            canonical['participantes'].append({field:record[field] for field in ('documento','nombre_completo') if record.get(field)})
        else: canonical['participantes'].append(record)


def validate_against_master(database_path: str, tenant_id: int, canonical: dict) -> dict:
    participants = list(canonical.get('participantes') or [])
    results: list[dict] = []
    if canonical.get('tipo_documento') in {'LISTADO_ASISTENCIA','RAM'} and not participants:
        results.append({'ruta_canonica':'participantes','regla':'PARTICIPANTES_REQUERIDOS','nivel':'CRITICO','estado':'ERROR','mensaje':'No se identificaron filas de participantes. Revisa el mapeo antes de aprobar.','esperado':None,'evidencia':{}})
    try:
        conn = connect(database_path)
        rows = conn.execute("""SELECT documento,nombre_completo,unidad_servicio,estado,activo FROM master_ninos WHERE COALESCE(fundacion_id,1)=?""", (tenant_id,)).fetchall()
        conn.close()
    except Exception:
        return {'semaforo':'GRIS','errores_criticos':0,'advertencias':1,'coincidencias':0,'total':len(participants),'resultados':[{'ruta_canonica':'participantes','regla':'BASE_MAESTRA_DISPONIBLE','nivel':'ADVERTENCIA','estado':'PENDIENTE','mensaje':'La Base Maestra no está disponible para validar este documento.','esperado':None,'evidencia':{}}]}
    master = {_document_key(row['documento']):dict(row) for row in rows if _document_key(row['documento'])}
    seen: set[str] = set(); matches = 0
    document_unit = normalize((canonical.get('unidad_servicio') or {}).get('nombre'))
    for index, participant in enumerate(participants):
        path=f'participantes.{index}'; document=_document_key(participant.get('documento'))
        if not document:
            results.append({'ruta_canonica':f'{path}.documento','regla':'DOCUMENTO_OBLIGATORIO','nivel':'CRITICO','estado':'ERROR','mensaje':'No se encontró documento o NUI.','esperado':None,'evidencia':{'indice':index}}); continue
        if document in seen:
            results.append({'ruta_canonica':f'{path}.documento','regla':'DUPLICADO_EN_DOCUMENTO','nivel':'CRITICO','estado':'ERROR','mensaje':'El participante está repetido en el documento cargado.','esperado':None,'evidencia':{'indice':index}})
        seen.add(document); found=master.get(document)
        if not found:
            results.append({'ruta_canonica':f'{path}.documento','regla':'EXISTE_BASE_MAESTRA','nivel':'CRITICO','estado':'ERROR','mensaje':'El documento no existe en la Base Maestra de esta fundación.','esperado':None,'evidencia':{'indice':index}}); continue
        matches += 1; participant['base_maestra_id_documento']=document; participant['validado_base_maestra']=True
        if not bool(found.get('activo')) or normalize(found.get('estado')) in {'retirado','inactivo'}:
            results.append({'ruta_canonica':path,'regla':'PARTICIPANTE_ACTIVO','nivel':'CRITICO','estado':'ERROR','mensaje':'El participante figura retirado o inactivo en Base Maestra.','esperado':{'estado':found.get('estado')},'evidencia':{'indice':index}})
        source_name=normalize(participant.get('nombre_completo')); expected_name=normalize(found.get('nombre_completo'))
        similarity=SequenceMatcher(None,source_name,expected_name).ratio() if source_name and expected_name else 0
        if similarity < .78:
            results.append({'ruta_canonica':f'{path}.nombre_completo','regla':'NOMBRE_COINCIDE_DOCUMENTO','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':'El nombre no coincide suficientemente con el documento en Base Maestra.','esperado':{'nombre_completo':found.get('nombre_completo')},'evidencia':{'similitud':round(similarity,3),'indice':index}})
        source_unit=normalize(participant.get('unidad') or document_unit); expected_unit=normalize(found.get('unidad_servicio'))
        if source_unit and expected_unit and source_unit != expected_unit:
            results.append({'ruta_canonica':f'{path}.unidad','regla':'UNIDAD_COINCIDE','nivel':'CRITICO','estado':'ERROR','mensaje':'El participante pertenece a otra UDS/UCA en Base Maestra.','esperado':{'unidad_servicio':found.get('unidad_servicio')},'evidencia':{'indice':index}})
    critical=sum(1 for item in results if item['nivel']=='CRITICO'); warnings=sum(1 for item in results if item['nivel']=='ADVERTENCIA')
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':matches,'total':len(participants),'resultados':results}


def validate_schedule(canonical: dict) -> dict:
    activities=list(canonical.get('actividades') or []); results=[]; seen=set()
    if not activities: results.append({'ruta_canonica':'actividades','regla':'ACTIVIDADES_REQUERIDAS','nivel':'CRITICO','estado':'ERROR','mensaje':'No se identificaron actividades en el cronograma.','esperado':None,'evidencia':{}})
    for index,activity in enumerate(activities):
        path=f'actividades.{index}'; date=activity.get('fecha') or activity.get('fecha_limite'); title=normalize(activity.get('actividad'))
        if not date: results.append({'ruta_canonica':f'{path}.fecha','regla':'FECHA_REQUERIDA','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':'La actividad no tiene una fecha válida; no se inventó ninguna.','esperado':None,'evidencia':{'indice':index}})
        key=(date,title)
        if key in seen: results.append({'ruta_canonica':path,'regla':'ACTIVIDAD_DUPLICADA','nivel':'CRITICO','estado':'ERROR','mensaje':'La actividad está repetida para la misma fecha.','esperado':None,'evidencia':{'indice':index}})
        seen.add(key)
        if not activity.get('responsable'): results.append({'ruta_canonica':f'{path}.responsable','regla':'RESPONSABLE_RECOMENDADO','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':'No se identificó responsable para la actividad.','esperado':None,'evidencia':{'indice':index}})
    critical=sum(item['nivel']=='CRITICO' for item in results); warnings=sum(item['nivel']=='ADVERTENCIA' for item in results)
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':len(activities),'total':len(activities),'resultados':results}


def validate_nutrition(database_path: str, tenant_id: int, canonical: dict) -> dict:
    records=list(canonical.get('valoraciones') or []); results=[]; matches=0; seen=set()
    try:
        conn=connect(database_path); rows=conn.execute("SELECT documento,nombre_completo,unidad_servicio,estado,activo FROM master_ninos WHERE COALESCE(fundacion_id,1)=?",(tenant_id,)).fetchall(); conn.close(); master={_document_key(row['documento']):dict(row) for row in rows}
    except Exception: master={}
    if not records: results.append({'ruta_canonica':'valoraciones','regla':'VALORACIONES_REQUERIDAS','nivel':'CRITICO','estado':'ERROR','mensaje':'No se identificaron filas de peso y talla.','esperado':None,'evidencia':{}})
    for index,record in enumerate(records):
        path=f'valoraciones.{index}'; document=_document_key(record.get('documento')); date=record.get('fecha'); key=(document,date)
        if not document: results.append({'ruta_canonica':f'{path}.documento','regla':'DOCUMENTO_OBLIGATORIO','nivel':'CRITICO','estado':'ERROR','mensaje':'La valoración no tiene documento o NUI.','esperado':None,'evidencia':{'indice':index}})
        elif document not in master: results.append({'ruta_canonica':f'{path}.documento','regla':'EXISTE_BASE_MAESTRA','nivel':'CRITICO','estado':'ERROR','mensaje':'El documento no existe en la Base Maestra de esta fundación.','esperado':None,'evidencia':{'indice':index}})
        else: matches+=1; record['validado_base_maestra']=True
        if key in seen: results.append({'ruta_canonica':path,'regla':'VALORACION_DUPLICADA','nivel':'CRITICO','estado':'ERROR','mensaje':'La valoración está repetida para el participante y fecha.','esperado':None,'evidencia':{'indice':index}})
        seen.add(key)
        if not date: results.append({'ruta_canonica':f'{path}.fecha','regla':'FECHA_REQUERIDA','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':'No se identificó una fecha válida para la valoración.','esperado':None,'evidencia':{'indice':index}})
        weight=record.get('peso_kg'); height=record.get('talla_cm'); arm=record.get('perimetro_braquial_cm')
        if weight is None or not 1<=weight<=120: results.append({'ruta_canonica':f'{path}.peso_kg','regla':'RANGO_PESO','nivel':'CRITICO','estado':'ERROR','mensaje':'El peso debe ser numérico y estar entre 1 y 120 kg.','esperado':{'min':1,'max':120},'evidencia':{'indice':index}})
        if height is None or not 30<=height<=220: results.append({'ruta_canonica':f'{path}.talla_cm','regla':'RANGO_TALLA','nivel':'CRITICO','estado':'ERROR','mensaje':'La talla debe ser numérica y estar entre 30 y 220 cm.','esperado':{'min':30,'max':220},'evidencia':{'indice':index}})
        if arm is not None and not 5<=arm<=60: results.append({'ruta_canonica':f'{path}.perimetro_braquial_cm','regla':'RANGO_PERIMETRO_BRAQUIAL','nivel':'CRITICO','estado':'ERROR','mensaje':'El perímetro braquial está fuera del rango técnico configurado.','esperado':{'min':5,'max':60},'evidencia':{'indice':index}})
    critical=sum(item['nivel']=='CRITICO' for item in results); warnings=sum(item['nivel']=='ADVERTENCIA' for item in results)
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':matches,'total':len(records),'resultados':results}


def validate_planning(canonical: dict) -> dict:
    planning=canonical.get('planeacion') or {}; results=[]
    for field,label in (('objetivo','objetivo o intencionalidad pedagógica'),('actividad','actividad o experiencia pedagógica')):
        if not planning.get(field): results.append({'ruta_canonica':f'planeacion.{field}','regla':f'{field.upper()}_REQUERIDO','nivel':'CRITICO','estado':'ERROR','mensaje':f'No se identificó {label}.','esperado':None,'evidencia':{}})
    for field,label in (('fecha_programada','fecha programada'),('responsable','responsable'),('unidad','UDS/UCA')):
        if not planning.get(field): results.append({'ruta_canonica':f'planeacion.{field}','regla':f'{field.upper()}_RECOMENDADO','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':f'No se identificó {label}; el sistema no completó el dato por suposición.','esperado':None,'evidencia':{}})
    critical=sum(item['nivel']=='CRITICO' for item in results); warnings=sum(item['nivel']=='ADVERTENCIA' for item in results)
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':len(planning),'total':len(PLANNING_LABEL_ALIASES),'resultados':results}


def _validate_labeled(data: dict, key: str, required: tuple[tuple[str,str],...], recommended: tuple[tuple[str,str],...], total_fields: int) -> dict:
    results=[]
    for field,label in required:
        if not data.get(field): results.append({'ruta_canonica':f'{key}.{field}','regla':f'{field.upper()}_REQUERIDO','nivel':'CRITICO','estado':'ERROR','mensaje':f'No se identificó {label}.','esperado':None,'evidencia':{}})
    for field,label in recommended:
        if not data.get(field): results.append({'ruta_canonica':f'{key}.{field}','regla':f'{field.upper()}_RECOMENDADO','nivel':'ADVERTENCIA','estado':'REVISAR','mensaje':f'No se identificó {label}; no se completó por suposición.','esperado':None,'evidencia':{}})
    critical=sum(item['nivel']=='CRITICO' for item in results); warnings=sum(item['nivel']=='ADVERTENCIA' for item in results)
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':len(data),'total':total_fields,'resultados':results}


def validate_minutes(canonical: dict) -> dict:
    return _validate_labeled(canonical.get('acta') or {},'acta',(('fecha','fecha del acta'),('tema','tema u objetivo'),('desarrollo','desarrollo de la reunión'),('compromisos','compromisos o acuerdos')),(('responsable','responsable'),('asistentes','asistentes'),('firma_referencia','presencia de firmas')),len(MINUTES_LABEL_ALIASES))


def validate_report(canonical: dict) -> dict:
    return _validate_labeled(canonical.get('informe') or {},'informe',(('objetivo','objetivo del informe'),('actividades','actividades realizadas'),('resultados','resultados'),('conclusiones','conclusiones')),(('periodo','periodo'),('responsable','responsable')),len(REPORT_LABEL_ALIASES))


def validate_rpp(canonical: dict) -> dict:
    return _validate_labeled(canonical.get('rpp') or {},'rpp',(('fecha','fecha del servicio'),('unidad','UDS/UCA'),('tiempo_comida','tiempo de comida'),('preparacion','preparación o menú'),('porciones','número de porciones')),(('minuta_patron','minuta patrón'),('responsable','responsable de la preparación')),len(RPP_LABEL_ALIASES))


def validate_bienestarina(database_path: str, tenant_id: int, canonical: dict) -> dict:
    validation=validate_against_master(database_path,tenant_id,canonical); results=list(validation.get('resultados') or [])
    deliveries=list(canonical.get('entregas') or [])
    if not deliveries: results.append({'ruta_canonica':'entregas','regla':'ENTREGAS_REQUERIDAS','nivel':'CRITICO','estado':'ERROR','mensaje':'No se identificaron entregas de Bienestarina.','esperado':None,'evidencia':{}})
    for index,delivery in enumerate(deliveries):
        for field,label in (('fecha_entrega','fecha de entrega'),('lote','lote'),('cantidad','cantidad entregada')):
            if not delivery.get(field): results.append({'ruta_canonica':f'entregas.{index}.{field}','regla':f'{field.upper()}_REQUERIDO','nivel':'CRITICO','estado':'ERROR','mensaje':f'No se identificó {label}.','esperado':None,'evidencia':{'indice':index}})
        quantity=_numeric_measure(delivery.get('cantidad'))
        if delivery.get('cantidad') and (quantity is None or quantity<=0): results.append({'ruta_canonica':f'entregas.{index}.cantidad','regla':'CANTIDAD_POSITIVA','nivel':'CRITICO','estado':'ERROR','mensaje':'La cantidad entregada debe ser numérica y mayor que cero.','esperado':{'min_exclusivo':0},'evidencia':{'indice':index}})
    critical=sum(item['nivel']=='CRITICO' for item in results); warnings=sum(item['nivel']=='ADVERTENCIA' for item in results)
    return {'semaforo':'ROJO' if critical else ('AMARILLO' if warnings else 'VERDE'),'errores_criticos':critical,'advertencias':warnings,'coincidencias':validation.get('coincidencias',0),'total':len(deliveries),'resultados':results}


def validate_canonical(database_path: str, tenant_id: int, canonical: dict) -> dict:
    if canonical.get('tipo_documento')=='CRONOGRAMA': return validate_schedule(canonical)
    if canonical.get('tipo_documento')=='PESO_TALLA': return validate_nutrition(database_path,tenant_id,canonical)
    if canonical.get('tipo_documento')=='PLANEACION_PEDAGOGICA': return validate_planning(canonical)
    if canonical.get('tipo_documento')=='ACTA': return validate_minutes(canonical)
    if canonical.get('tipo_documento')=='INFORME': return validate_report(canonical)
    if canonical.get('tipo_documento')=='RPP': return validate_rpp(canonical)
    if canonical.get('tipo_documento')=='BIENESTARINA': return validate_bienestarina(database_path,tenant_id,canonical)
    return validate_against_master(database_path,tenant_id,canonical)


def attendance_official_payload(document: dict) -> tuple[list[dict], dict]:
    if not document or document.get('estado') != 'APROBADO':
        raise ValueError('Aprueba el documento antes de generar el listado oficial.')
    if document.get('tipo_documento') != 'LISTADO_ASISTENCIA':
        raise ValueError('La generación oficial solo está disponible para listados de asistencia.')
    canonical = document.get('resultado_canonico') or {}
    participants = list(canonical.get('participantes') or [])
    if not participants:
        raise ValueError('El documento aprobado no contiene participantes para generar.')
    unit = (canonical.get('unidad_servicio') or {}).get('nombre') or ''
    period = canonical.get('periodo') or {}
    users = []
    for participant in participants:
        users.append({
            'nombre_completo': participant.get('nombre_completo') or '',
            'documento': participant.get('documento') or participant.get('nui') or '',
            'tipo_documento': participant.get('tipo_documento') or '',
            'unidad': participant.get('unidad') or unit,
            'telefono': participant.get('telefono') or '',
        })
    return users, {'unidad':unit,'fecha':period.get('fecha'),'mes':period.get('mes'),'anio':period.get('anio'),'documento_id':document.get('id')}


def public_document(row: Any) -> dict:
    item = dict(row)
    for source, target, default in (
        ('resultado_canonico_json', 'resultado_canonico', {}),
        ('validaciones_json', 'validaciones', []),
    ):
        try:
            item[target] = json.loads(item.get(source) or json.dumps(default))
        except Exception:
            item[target] = default
        item.pop(source, None)
    item.pop('resultado_bruto_json', None)
    item.pop('ruta_privada', None)
    return item
