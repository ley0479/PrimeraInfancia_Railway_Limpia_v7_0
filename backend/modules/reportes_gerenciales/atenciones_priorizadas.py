from __future__ import annotations

import hashlib
import json
import os
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from .services import ReportesGerencialesService, normalizar_texto, periodo_key


ATENCIONES = (
    ("EDUCACION_INICIAL", "Educación Inicial", "Cuéntame, calendario, planeaciones y asistencia", "SEMIAUTOMATICA"),
    ("DOCUMENTO_IDENTIDAD", "Documento de Identidad", "Base Maestra / Cuéntame", "AUTOMATICA"),
    ("TALENTO_HUMANO", "Talento Humano Cualificado", "Talento Humano", "AUTOMATICA"),
    ("MATERIALES_LITERARIOS", "Materiales literarios especializados", "Lista de chequeo mensual por UDS", "SEMIAUTOMATICA"),
    ("ESTADO_NUTRICIONAL", "Seguimiento del estado nutricional", "Salud y Nutrición", "AUTOMATICA"),
    ("FORMACION_FAMILIAS", "Formación a familias y cuidadores", "Calendario, actividades, listados y evidencias", "SEMIAUTOMATICA"),
    ("CRECIMIENTO_DESARROLLO", "Carné de crecimiento y desarrollo", "Salud / lectura documental confirmada", "SEMIAUTOMATICA"),
    ("AFILIACION_SALUD", "Afiliación vigente a salud", "Base de afiliación / soporte confirmado", "SEMIAUTOMATICA"),
    ("VACUNACION", "Vacunación al día", "Registro de vacunación / soporte confirmado", "SEMIAUTOMATICA"),
)


class AtencionesPriorizadasService:
    def __init__(self, base: ReportesGerencialesService):
        self.base = base

    @staticmethod
    def _doc_type(row: dict[str, Any]) -> str:
        raw = normalizar_texto(row.get("tipo_documento") or row.get("tipo_documento_beneficiario") or row.get("tipo_doc"))
        if "registro" in raw or raw == "rc": return "RC"
        if "tarjeta" in raw or raw == "ti": return "TI"
        if "cedula extranjeria" in raw or raw == "ce": return "CE"
        if "cedula" in raw or raw == "cc": return "CC"
        return "SD"

    def consolidar(self, fundacion_id: int, mes: int, anio: int, cobertura: int = 0) -> dict[str, Any]:
        base_data = self.base.recopilar_datos(mes, anio, fundacion_id)
        personas = self.base.get_beneficiarios(fundacion_id)
        talento = self.base.get_talento(fundacion_id)
        total = len(personas)
        cobertura = int(cobertura or total)
        documentos: dict[str, int] = {"CC": 0, "CE": 0, "RC": 0, "SD": 0, "TI": 0}
        for row in personas:
            documentos[self._doc_type(row)] = documentos.get(self._doc_type(row), 0) + 1

        grupos = base_data.get("distribucion_grupos") or {}
        total_rpp = sum(int(v or 0) for v in grupos.values())
        salud = base_data.get("salud") or {}
        valorados = int(salud.get("valorados") or 0)
        entregables = base_data.get("entregables") or {}
        formacion = sum(1 for item in entregables.get("items") or [] if any(k in normalizar_texto(item.get("titulo") or item.get("tipo")) for k in ("famil", "cuidador", "encuentro")))

        automaticos = {
            "EDUCACION_INICIAL": (total, cobertura, {"usuarios_activos": total, "uds": len(base_data.get("distribucion_unidades") or {})}),
            "DOCUMENTO_IDENTIDAD": (total - documentos.get("SD", 0), total, {"tipos": documentos, "total": sum(documentos.values())}),
            "TALENTO_HUMANO": (len(talento), len(talento), {"total": len(talento), "cargos": base_data.get("distribucion_talento") or {}}),
            "ESTADO_NUTRICIONAL": (valorados, total, {"valorados": valorados, "criticos": salud.get("criticos", 0), "diagnosticos": salud.get("diagnosticos") or {}}),
            "FORMACION_FAMILIAS": (formacion, 1, {"encuentros_detectados": formacion}),
        }
        resultados = []
        for codigo, nombre, fuente, automatizacion in ATENCIONES:
            numerador, denominador, datos = automaticos.get(codigo, (0, total, {}))
            porcentaje = round((numerador * 100 / denominador), 2) if denominador else 0
            completo = denominador > 0 and numerador >= denominador
            estado = "COMPLETO" if completo else ("CON_ALERTAS" if numerador else "PENDIENTE")
            resultados.append({
                "codigo": codigo, "nombre": nombre, "fuente": fuente,
                "automatizacion": automatizacion, "numerador": numerador,
                "denominador": denominador, "porcentaje": porcentaje,
                "estado": estado, "datos": datos,
            })

        hallazgos = []
        if sum(documentos.values()) != total:
            hallazgos.append({"codigo": "DOCUMENTOS_TOTAL_DIFIERE", "atencion": "DOCUMENTO_IDENTIDAD", "nivel": "ERROR", "mensaje": "Los tipos documentales no suman el total de beneficiarios."})
        if total_rpp != total:
            hallazgos.append({"codigo": "RPP_TOTAL_DIFIERE", "atencion": "EDUCACION_INICIAL", "nivel": "ERROR", "mensaje": f"Las categorías RPP suman {total_rpp} y la Base Maestra contiene {total} participantes.", "diferencia": total - total_rpp})
        return {
            "periodo": periodo_key(anio, mes), "mes": mes, "anio": anio,
            "fundacion_id": fundacion_id, "fundacion": base_data.get("fundacion_nombre"),
            "cobertura_contratada": cobertura, "beneficiarios": total,
            "rpp_total": total_rpp, "resultados": resultados, "hallazgos": hallazgos,
            "avance": {"completas": sum(1 for x in resultados if x["estado"] == "COMPLETO"), "total": 9},
            "generado_en": datetime.now().isoformat(timespec="seconds"),
        }

    def guardar_borrador(self, payload: dict[str, Any], user: dict[str, Any]) -> dict[str, Any]:
        fid = int(user.get("fundacion_id") or 1)
        mes, anio = int(payload.get("mes")), int(payload.get("anio"))
        data = self.consolidar(fid, mes, anio, int(payload.get("cobertura_contratada") or 0))
        conn = self.base.connect(); cur = conn.cursor(); now = datetime.now().isoformat(timespec="seconds")
        contrato = str(payload.get("contrato") or "").strip()
        row = cur.execute("SELECT id FROM rg9_informes WHERE fundacion_id=? AND contrato=? AND periodo=? AND version=1", (fid, contrato, data["periodo"])).fetchone()
        if row:
            informe_id = int(row["id"])
            cur.execute("UPDATE rg9_informes SET fecha_corte=?,cobertura_contratada=?,modalidad=?,actualizado_en=? WHERE id=? AND fundacion_id=?", (payload.get("fecha_corte"), data["cobertura_contratada"], payload.get("modalidad"), now, informe_id, fid))
        else:
            cur.execute("INSERT INTO rg9_informes(fundacion_id,contrato,periodo,mes,anio,fecha_corte,cobertura_contratada,modalidad,estado,version,responsable_id,creado_en) VALUES(?,?,?,?,?,?,?,?, 'BORRADOR',1,?,?)", (fid, contrato, data["periodo"], mes, anio, payload.get("fecha_corte"), data["cobertura_contratada"], payload.get("modalidad"), user.get("id"), now))
            informe_id = int(cur.lastrowid)
        for item in data["resultados"]:
            cur.execute("INSERT INTO rg9_resultados(informe_id,atencion_codigo,numerador,denominador,porcentaje,estado,fuente,fecha_actualizacion,datos_json) VALUES(?,?,?,?,?,?,?,?,?) ON CONFLICT(informe_id,atencion_codigo) DO UPDATE SET numerador=excluded.numerador,denominador=excluded.denominador,porcentaje=excluded.porcentaje,estado=excluded.estado,fuente=excluded.fuente,fecha_actualizacion=excluded.fecha_actualizacion,datos_json=excluded.datos_json", (informe_id,item["codigo"],item["numerador"],item["denominador"],item["porcentaje"],item["estado"],item["fuente"],now,json.dumps(item["datos"],ensure_ascii=False)))
        cur.execute("DELETE FROM rg9_hallazgos WHERE informe_id=?", (informe_id,))
        for h in data["hallazgos"]:
            cur.execute("INSERT INTO rg9_hallazgos(informe_id,atencion_codigo,codigo,nivel,mensaje,estado,creado_en) VALUES(?,?,?,?,?,'ABIERTO',?)", (informe_id,h.get("atencion"),h["codigo"],h["nivel"],h["mensaje"],now))
        conn.commit(); conn.close()
        return {"id": informe_id, "estado": "BORRADOR", **data}

    def aprobar(self, informe_id: int, user: dict[str, Any]) -> dict[str, Any]:
        fid = int(user.get("fundacion_id") or 1); conn = self.base.connect(); cur = conn.cursor()
        report = cur.execute("SELECT * FROM rg9_informes WHERE id=? AND fundacion_id=?", (informe_id, fid)).fetchone()
        if not report: conn.close(); raise LookupError("Informe no encontrado para esta fundación.")
        abiertos = cur.execute("SELECT COUNT(*) c FROM rg9_hallazgos WHERE informe_id=? AND nivel='ERROR' AND estado='ABIERTO'", (informe_id,)).fetchone()["c"]
        pendientes = cur.execute("SELECT COUNT(*) c FROM rg9_resultados WHERE informe_id=? AND estado!='COMPLETO'", (informe_id,)).fetchone()["c"]
        if abiertos or pendientes:
            conn.close(); raise ValueError(f"No es posible aprobar: {abiertos} inconsistencia(s) y {pendientes} atención(es) incompletas.")
        results = [dict(x) for x in cur.execute("SELECT * FROM rg9_resultados WHERE informe_id=? ORDER BY id", (informe_id,)).fetchall()]
        snapshot = {"informe": dict(report), "resultados": results}
        raw = json.dumps(snapshot, ensure_ascii=False, sort_keys=True, default=str)
        digest = hashlib.sha256(raw.encode("utf-8")).hexdigest(); now = datetime.now().isoformat(timespec="seconds")
        cur.execute("INSERT INTO rg9_snapshots(informe_id,fundacion_id,version,datos_json,hash_sha256,creado_en) VALUES(?,?,?,?,?,?)", (informe_id,fid,int(report["version"]),raw,digest,now))
        cur.execute("UPDATE rg9_informes SET estado='APROBADO',aprobado_en=?,actualizado_en=? WHERE id=? AND fundacion_id=?", (now,now,informe_id,fid))
        conn.commit(); conn.close(); return {"id": informe_id, "estado": "APROBADO", "snapshot_hash": digest}

    def detalle(self, informe_id: int, fundacion_id: int) -> dict[str, Any]:
        conn = self.base.connect()
        report = conn.execute("SELECT * FROM rg9_informes WHERE id=? AND fundacion_id=?", (informe_id, fundacion_id)).fetchone()
        if not report: conn.close(); raise LookupError("Informe no encontrado para esta fundación.")
        resultados = [dict(x) for x in conn.execute("SELECT * FROM rg9_resultados WHERE informe_id=? ORDER BY id", (informe_id,)).fetchall()]
        hallazgos = [dict(x) for x in conn.execute("SELECT * FROM rg9_hallazgos WHERE informe_id=? ORDER BY id", (informe_id,)).fetchall()]
        evidencias = [dict(x) for x in conn.execute("SELECT id,atencion_codigo,unidad,nombre_archivo,fecha_evidencia,responsable,estado_revision,creado_en FROM rg9_evidencias WHERE informe_id=? ORDER BY id", (informe_id,)).fetchall()]
        conn.close()
        for item in resultados:
            try: item["datos"] = json.loads(item.pop("datos_json") or "{}")
            except Exception: item["datos"] = {}
        return {"informe": dict(report), "resultados": resultados, "hallazgos": hallazgos, "evidencias": evidencias}

    def actualizar_atencion(self, informe_id: int, codigo: str, payload: dict[str, Any], user: dict[str, Any]) -> dict[str, Any]:
        fid = int(user.get("fundacion_id") or 1); codigo = str(codigo or "").upper()
        if codigo not in {x[0] for x in ATENCIONES}: raise ValueError("Atención no reconocida.")
        conn = self.base.connect(); report = conn.execute("SELECT estado FROM rg9_informes WHERE id=? AND fundacion_id=?", (informe_id,fid)).fetchone()
        if not report: conn.close(); raise LookupError("Informe no encontrado para esta fundación.")
        if report["estado"] in {"APROBADO", "CERRADO"}: conn.close(); raise ValueError("El informe aprobado está congelado.")
        numerador = max(0, int(payload.get("numerador") or 0)); denominador = max(0, int(payload.get("denominador") or 0))
        if numerador > denominador and denominador: conn.close(); raise ValueError("El numerador no puede superar el denominador.")
        porcentaje = round(numerador * 100 / denominador, 2) if denominador else 0
        estado = "COMPLETO" if denominador and numerador >= denominador else ("CON_ALERTAS" if numerador else "PENDIENTE")
        now = datetime.now().isoformat(timespec="seconds")
        conn.execute("UPDATE rg9_resultados SET numerador=?,denominador=?,porcentaje=?,estado=?,observacion=?,responsable=?,datos_json=?,fecha_actualizacion=? WHERE informe_id=? AND atencion_codigo=?", (numerador,denominador,porcentaje,estado,payload.get("observacion"),payload.get("responsable"),json.dumps(payload.get("datos") or {},ensure_ascii=False),now,informe_id,codigo))
        conn.execute("UPDATE rg9_informes SET actualizado_en=? WHERE id=? AND fundacion_id=?", (now,informe_id,fid)); conn.commit(); conn.close()
        return {"codigo": codigo, "numerador": numerador, "denominador": denominador, "porcentaje": porcentaje, "estado": estado}

    def _export_dir(self, report: dict[str, Any]) -> Path:
        folder = self.base.reportes_folder / "9_atenciones" / str(report["periodo"])
        folder.mkdir(parents=True, exist_ok=True); return folder

    def guardar_evidencia(self, informe_id: int, codigo: str, upload: Any, payload: dict[str, Any], user: dict[str, Any]) -> dict[str, Any]:
        from werkzeug.utils import secure_filename
        fid=int(user.get("fundacion_id") or 1); codigo=str(codigo or "").upper()
        if codigo not in {x[0] for x in ATENCIONES}: raise ValueError("Atención no reconocida.")
        conn=self.base.connect(); report=conn.execute("SELECT * FROM rg9_informes WHERE id=? AND fundacion_id=?",(informe_id,fid)).fetchone()
        if not report: conn.close(); raise LookupError("Informe no encontrado para esta fundación.")
        if report["estado"] in {"APROBADO","CERRADO"}: conn.close(); raise ValueError("El informe aprobado está congelado.")
        name=secure_filename(getattr(upload,"filename","") or "")
        ext=Path(name).suffix.lower()
        if ext not in {".pdf",".xlsx",".xls",".docx",".png",".jpg",".jpeg"}: conn.close(); raise ValueError("Formato de evidencia no permitido.")
        folder=self._export_dir(dict(report))/"evidencias"/codigo; folder.mkdir(parents=True,exist_ok=True)
        stored=f"{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}_{name}"; path=folder/stored; upload.save(path)
        if not path.is_file() or path.stat().st_size<=0: conn.close(); raise ValueError("La evidencia llegó vacía.")
        now=datetime.now().isoformat(timespec="seconds")
        cur=conn.cursor(); cur.execute("INSERT INTO rg9_evidencias(informe_id,atencion_codigo,unidad,nombre_archivo,ruta_archivo,fecha_evidencia,responsable,estado_revision,creado_en) VALUES(?,?,?,?,?,?,?,?,?)",(informe_id,codigo,payload.get("unidad"),name,str(path),payload.get("fecha_evidencia"),payload.get("responsable"),"PENDIENTE",now)); evidence_id=int(cur.lastrowid); conn.commit(); conn.close()
        return {"id":evidence_id,"atencion_codigo":codigo,"nombre_archivo":name,"estado_revision":"PENDIENTE"}

    def guardar_plantilla_pptx(self, upload: Any, payload: dict[str, Any], user: dict[str, Any]) -> dict[str, Any]:
        from werkzeug.utils import secure_filename
        from pptx import Presentation
        fid=int(user.get("fundacion_id") or 1); name=secure_filename(getattr(upload,"filename","") or "")
        if Path(name).suffix.lower() != ".pptx": raise ValueError("La plantilla oficial debe ser un archivo .pptx.")
        folder=self.base.reportes_folder/"9_atenciones"/"plantillas"/str(fid); folder.mkdir(parents=True,exist_ok=True)
        stored=f"plantilla_9_atenciones_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.pptx"; path=folder/stored; upload.save(path)
        if not path.is_file() or path.stat().st_size<=0: raise ValueError("La plantilla llegó vacía.")
        try: presentation=Presentation(str(path)); slides=len(presentation.slides)
        except Exception as exc: path.unlink(missing_ok=True); raise ValueError("El PowerPoint no es una plantilla PPTX válida.") from exc
        digest=hashlib.sha256(path.read_bytes()).hexdigest(); conn=self.base.connect(); now=datetime.now().isoformat(timespec="seconds")
        conn.execute("UPDATE rg9_plantillas_pptx SET estado='HISTORICA' WHERE fundacion_id=? AND estado='ACTIVA'",(fid,))
        cur=conn.cursor(); cur.execute("INSERT INTO rg9_plantillas_pptx(fundacion_id,nombre_original,ruta_archivo,version,fecha_vigencia,estado,hash_sha256,cargado_por,creado_en) VALUES(?,?,?,?,?,'ACTIVA',?,?,?)",(fid,name,str(path),payload.get("version") or "1",payload.get("fecha_vigencia"),digest,user.get("id"),now)); template_id=int(cur.lastrowid); conn.commit(); conn.close()
        return {"id":template_id,"nombre_original":name,"version":payload.get("version") or "1","estado":"ACTIVA","hash_sha256":digest,"diapositivas":slides}

    def plantilla_pptx_activa(self, fundacion_id: int) -> dict[str, Any] | None:
        conn=self.base.connect(); row=conn.execute("SELECT id,nombre_original,ruta_archivo,version,fecha_vigencia,estado,hash_sha256,creado_en FROM rg9_plantillas_pptx WHERE fundacion_id=? AND estado='ACTIVA' ORDER BY id DESC LIMIT 1",(fundacion_id,)).fetchone(); conn.close()
        return dict(row) if row else None

    def generar_exportaciones(self, informe_id: int, fundacion_id: int) -> dict[str, str]:
        detail = self.detalle(informe_id, fundacion_id); report = detail["informe"]; results = detail["resultados"]
        folder = self._export_dir(report); base_name = f"9_ATENCIONES_{report['periodo']}_V{report['version']}"
        xlsx = folder / f"{base_name}_ANEXO.xlsx"; pptx = folder / f"{base_name}.pptx"; pdf = folder / f"{base_name}.pdf"; zpath = folder / f"{base_name}_EVIDENCIAS.zip"
        from openpyxl import Workbook
        wb = Workbook(); ws = wb.active; ws.title = "Resumen"
        ws.append(["Atención", "Numerador", "Denominador", "Porcentaje", "Estado", "Fuente", "Observación"])
        names = {x[0]: x[1] for x in ATENCIONES}
        for x in results: ws.append([names.get(x["atencion_codigo"],x["atencion_codigo"]),x["numerador"],x["denominador"],x["porcentaje"],x["estado"],x["fuente"],x.get("observacion")])
        wb.save(xlsx)
        from pptx import Presentation
        template=self.plantilla_pptx_activa(fundacion_id)
        template_path=template.get("ruta_archivo") if template else None
        prs = Presentation(template_path) if template_path and Path(template_path).is_file() else Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = "Informe mensual\n9 Atenciones Priorizadas"; slide.placeholders[1].text = f"Periodo: {report['periodo']}\nContrato: {report.get('contrato') or 'No configurado'}"
        for x in results:
            slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = names.get(x["atencion_codigo"],x["atencion_codigo"]); slide.placeholders[1].text = f"Resultado: {x['numerador']} / {x['denominador']} ({x['porcentaje']} %)\nEstado: {x['estado']}\nFuente: {x['fuente']}\nObservación: {x.get('observacion') or 'Sin observación'}"
        prs.save(pptx)
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        styles=getSampleStyleSheet(); story=[Paragraph("Informe mensual - 9 Atenciones Priorizadas",styles["Title"]),Paragraph(f"Periodo {report['periodo']}",styles["Heading2"]),Spacer(1,12)]
        for x in results: story.extend([Paragraph(names.get(x["atencion_codigo"],x["atencion_codigo"]),styles["Heading2"]),Paragraph(f"{x['numerador']} de {x['denominador']} ({x['porcentaje']} %) - {x['estado']}",styles["BodyText"]),Spacer(1,8)])
        SimpleDocTemplate(str(pdf),pagesize=letter).build(story)
        with zipfile.ZipFile(zpath,"w",zipfile.ZIP_DEFLATED) as archive:
            conn=self.base.connect(); rows=conn.execute("SELECT ruta_archivo,nombre_archivo,atencion_codigo,unidad FROM rg9_evidencias WHERE informe_id=?",(informe_id,)).fetchall(); conn.close()
            for row in rows:
                path=Path(os.fspath(row["ruta_archivo"]));
                if path.is_file(): archive.write(path,arcname=f"{row['atencion_codigo']}/{row['unidad'] or 'GENERAL'}/{row['nombre_archivo']}")
        return {"pptx":str(pptx),"pdf":str(pdf),"xlsx":str(xlsx),"zip":str(zpath)}
